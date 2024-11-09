from flask import Flask, render_template, request, redirect, url_for, session, flash, send_file
from langchain_groq import ChatGroq
from langchain_core.prompts import PromptTemplate
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import re
import os
from flask_mysqldb import MySQL
import logging

app = Flask(__name__)
app.secret_key = os.getenv('FLASK_SECRET_KEY', 'your_secret_key')  # Use an environment variable for production
app.config['MYSQL_HOST'] = 'localhost'
app.config['MYSQL_USER'] = 'abinav'
app.config['MYSQL_PASSWORD'] = 'Abinav@2006'
app.config['MYSQL_DB'] = 'ppt_database'

mysql = MySQL(app)

logging.basicConfig(level=logging.INFO)  # Set up logging

@app.route('/')
def home():
    if 'username' in session:
        return redirect(url_for('dashboard'))
    return redirect(url_for('register'))

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        
        cur = mysql.connection.cursor()
        try:
            cur.execute("INSERT INTO users(username, password) VALUES(%s, %s)", (username, password))
            mysql.connection.commit()
            flash('Registration successful! Please log in.', 'success')
        except Exception as e:
            logging.error(f"Error registering user: {e}")
            flash('Registration failed. Please try again.', 'danger')
        finally:
            cur.close()
        
        return redirect(url_for('login'))

    return render_template('register.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        
        cur = mysql.connection.cursor()
        cur.execute("SELECT * FROM users WHERE username = %s AND password = %s", (username, password))
        user = cur.fetchone()
        cur.close()
        
        if user:
            session['username'] = username
            return redirect(url_for('dashboard'))
        else:
            flash('Invalid credentials. Please try again.', 'danger')
    
    return render_template('login.html')

@app.route('/dashboard')
def dashboard():
    if 'username' not in session:
        return redirect(url_for('login'))
    return render_template('dashboard.html', username=session['username'])

@app.route('/logout')
def logout():
    session.pop('username', None)
    flash('You have been logged out.', 'success')
    return redirect(url_for('login'))

@app.route('/update_password', methods=['GET', 'POST'])
def update_password():
    if 'username' not in session:
        return redirect(url_for('login'))
    
    if request.method == 'POST':
        new_password = request.form['new_password']
        cur = mysql.connection.cursor()
        try:
            cur.execute("UPDATE users SET password = %s WHERE username = %s", (new_password, session['username']))
            mysql.connection.commit()
            flash('Password updated successfully!', 'success')
        except Exception as e:
            logging.error(f"Error updating password: {e}")
            flash('Failed to update password. Please try again.', 'danger')
        finally:
            cur.close()
        return redirect(url_for('dashboard'))

    return render_template('update_password.html')
llm = ChatGroq(
    temperature=0.5,
    groq_api_key='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',
    model_name="llama-3.1-70b-versatile"
)

prompt = PromptTemplate.from_template(
    """
    Using the inputs provided:

Number of Slides: {n}
Topic: {topic}
Create a PowerPoint presentation with the following specifications:

Each slide should include relevant content gathered from trusted online sources, using the provided reference links to ensure accuracy and credibility.
Integrate high-quality, topic-related images and graphics sourced from the internet to enhance visual appeal and clarity.
Use a modern and professional design style, with consistent fonts set to Arial, harmonious color schemes, and visually engaging elements, such as shapes, gradients, or subtle animations.
Arrange the information and visuals to maintain a clean layout, making the content concise and easy to understand.
Include a final slide listing all reference sources for both content and images used throughout the presentation.
Aim for an informative and aesthetically pleasing presentation, with each slide showcasing a balanced mix of text and visuals in line with the specified topic and references.
Give a humanised content. Do not give any preamble.

Do not include any markdown langiage give only pure text

Format the output as follows:
Slide: number of the slide.
For the Title slide:
Title:
Subtitle:

For the other slides:
Title:
Content: Content: Keep the text on each slide brief and to the point and you should generate some examples and give it in a human like way. Depending on the nature of the topic, your approach will vary:
if {topic} is technical then give explanation of the {topic} in detail with 500 words paragraphs and it should have and give technical examples and explain them in detail and if it contains programming to generate a minimum of 5 code based examples on the slide topic and explain them in detail. each code example should have a seperate slide number.
if the {topic} is non technical then give explanations as a paragraph and in some slides give some points.
based on n{topic} 
No need of Images

For Fun, Engaging Points or Surprising Insights:
Occasionally, include "wow" moments—something that catches the audience's attention and adds excitement. These can be surprising facts, quick trivia, or interesting real-world applications that bring the topic to life.
Catchy Points: Use these engaging slides to reinforce a key concept, making the information memorable.

Do not give single point give as a 500 words

For the Conclusion or the Last Slide
Title: Conclusion
Content: Use this slide to wrap up the presentation by summarizing the key takeaways or lessons learned. Keep it brief and impactful. Consider using a call to action, a reflective question, or a final thought that encourages the audience to apply what they’ve learned or think about the subject in a new way. 
    """
)



@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate', methods=['POST'])
def generate_ppt():
    topic = request.form['topic']
    n = int(request.form['number_of_slides'])
    
    
    chain = prompt | llm
    result = chain.invoke(input={'topic': topic, 'n': n})
    
    content = result.content
    slides_content = parse_presentation_content(content)
    
    pptx_file = create_ppt(slides_content)
    
    return send_file(pptx_file, as_attachment=True)

def parse_presentation_content(content):
    slides = []
    slide_sections = content.strip().split("Slide")
    
    for section in slide_sections:
        if not section.strip():
            continue

        slide_number_match = re.search(r"\d+:", section)
        slide_number = int(slide_number_match.group(0).replace(":", "").strip()) if slide_number_match else None

        title_match = re.search(r"Title:\s(.*)", section)
        title = title_match.group(1).strip() if title_match else ""

        subtitle = ""
        if slide_number == 1:
            subtitle_match = re.search(r"Subtitle:\s(.*)", section)
            subtitle = subtitle_match.group(1).strip() if subtitle_match else ""

       # Extract content or bullet points
        content_match = re.search(r"Content:\s*(.*)", section, re.DOTALL)
        if content_match:
            content = content_match.group(1).strip()
            # Split content into bullet points based on bullets or lines
            bullet_points = [point.strip() for point in re.split(r"•|\n", content) if point.strip()]
        else:
            bullet_points = []

        # Append parsed slide data
        slides.append({
            "slide_number": slide_number,
            "title": title,
            "subtitle": subtitle,
            "bullet_points": bullet_points
        })

    return slides

def create_ppt(slides_content):
    prs = Presentation()
    
    for slide_data in slides_content:
        # Choose a layout for each slide type
        if slide_data["slide_number"] == 1:
            # Title Slide layout
            slide_layout = prs.slide_layouts[0]  # Title Slide layout
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = slide_data["title"]

            # Set modern style for title
            title_shape.text_frame.paragraphs[0].font.size = Pt(44)
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 128)

            if slide_data["subtitle"]:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = slide_data["subtitle"]
                subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
                subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)

        elif slide_data["title"].lower() == "conclusion":
            # Slide for Conclusion (without bullet points)
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = slide_data["title"]
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True

            content = slide.placeholders[1]
            content.text = slide_data["bullet_points"][0]  # Full text content for the conclusion
            content.text_frame.paragraphs[0].font.size = Pt(20)
            content.text_frame.paragraphs[0].font.color.rgb = RGBColor(60, 60, 60)

        else:
            # Regular slides with bullet points
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Title setup
            title = slide.shapes.title
            title.text = slide_data["title"]
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Dark Blue

            # Bullet points
            content = slide.placeholders[1]
            for point in slide_data["bullet_points"]:
                p = content.text_frame.add_paragraph()
                p.text = point
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(80, 80, 80)
                p.level = 0  # Bullet point level for uniformity

    pptx_file = "Modern_Presentation.pptx"
    prs.save(pptx_file)
    return pptx_file

if __name__ == '__main__':
    app.run(debug=True)
